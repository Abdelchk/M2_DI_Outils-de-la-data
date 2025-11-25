"""
Script pour créer une présentation PowerPoint à partir du markdown
Nécessite python-pptx : pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Erreur : python-pptx n'est pas installé.")
    print("Installez-le avec : pip install python-pptx")
    exit(1)

import re

def creer_presentation():
    """Crée une présentation PowerPoint depuis le markdown"""
    
    # Lire le fichier markdown
    with open('cours-complet-50-slides.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Créer une présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Couleurs
    couleur_titre = RGBColor(44, 62, 80)  # Bleu foncé
    couleur_texte = RGBColor(0, 0, 0)     # Noir
    couleur_accent = RGBColor(52, 152, 219)  # Bleu
    
    # Fonction pour créer un slide titre
    def slide_titre(titre, sous_titre=None):
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout titre
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = titre
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = couleur_titre
        
        if sous_titre:
            subtitle.text = sous_titre
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = couleur_accent
    
    # Fonction pour créer un slide de contenu
    def slide_contenu(titre, contenu):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Layout titre et contenu
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = titre
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = couleur_titre
        
        tf = content_shape.text_frame
        tf.word_wrap = True
        tf.clear()
        
        # Ajouter le contenu
        p = tf.paragraphs[0]
        p.text = contenu
        p.font.size = Pt(18)
        p.font.color.rgb = couleur_texte
        p.space_after = Pt(12)
    
    # Slide de titre
    slide_titre("OUTILS DE LA DATA", "Master 2 - Data Intelligence\nFormateur : Abid Hamza")
    
    # Parser le markdown - format avec SLIDE X
    slides_data = re.findall(r'^# SLIDE \d+ : (.+?)\n\n(.*?)(?=^# SLIDE|\Z)', content, re.MULTILINE | re.DOTALL)
    
    for titre, contenu_raw in slides_data:
        # Nettoyer le contenu
        contenu_lines = []
        for line in contenu_raw.split('\n'):
            line = line.strip()
            if line and not line.startswith('━') and not line.startswith('**'):
                # Nettoyer le markdown
                line = re.sub(r'^\*\*', '', line)
                line = re.sub(r'\*\*$', '', line)
                line = re.sub(r'^### ', '', line)
                line = re.sub(r'^## ', '', line)
                line = re.sub(r'^- ', '• ', line)
                line = re.sub(r'^\d+\. ', '', line)
                # Préserver les schémas ASCII
                if line.startswith('┌') or line.startswith('│') or line.startswith('└') or line.startswith(' '):
                    contenu_lines.append(line)
                elif line:
                    contenu_lines.append(line)
        
        contenu = '\n'.join(contenu_lines[:20])  # Limiter à 20 lignes par slide
        
        if contenu.strip():
            slide_contenu(titre, contenu)
    
    # Slide de conclusion
    slide_titre("MERCI POUR VOTRE ATTENTION", "Master 2 - Data Intelligence\nOutils de la Data")
    
    # Sauvegarder
    prs.save('cours-outils-data-complet.pptx')
    print("Presentation PowerPoint creee : cours-outils-data-complet.pptx")
    print(f"Nombre total de slides : {len(prs.slides)}")

if __name__ == '__main__':
    creer_presentation()

